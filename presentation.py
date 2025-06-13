from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor

# Skapa en ny presentation
prs = Presentation()

# Definiera en titelbild
def add_title_slide(title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

# Skapa vanliga innehållsbilder
def add_content_slide(title, bullet_points):
    slide_layout = prs.slide_layouts[1]  # Titel och innehåll
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    content = slide.placeholders[1]
    content.text = bullet_points[0]
    for point in bullet_points[1:]:
        content.text += f"\n{point}"

# Slide 1: Titel
add_title_slide(
    "Varför Python är ett bättre val än MATLAB",
    "Effektivare, billigare och mer framtidssäkert\n(Presentatörens namn – Datum)"
)

# Slide 2: Ekonomi & Tillgänglighet
add_content_slide("Python är gratis. MATLAB kostar.", [
    "• Python är helt öppen källkod – inga licenser, inga begränsningar.",
    "• MATLAB kräver dyra licenser + extra kostnader för toolboxar.",
    "• Python sparar företaget pengar – både på kort och lång sikt.",
    "🟢 Kostnadseffektivt\n🟢 Skalbart för hela teamet"
])

# Slide 3: Flexibilitet & Integration
add_content_slide("Python är mångsidigt och modernt", [
    "• Används för dataanalys, automation, AI/ML, API-integrationer, webbutveckling m.m.",
    "• Enkelt att koppla till databaser, REST API:er, molntjänster.",
    "• Fungerar bra i både Linux-, Windows- och molnmiljöer.",
    "🛠️ Fler användningsområden\n🔗 Bättre ekosystem"
])

# Slide 4: Community & Utveckling
add_content_slide("Python har ett levande ekosystem", [
    "• Stort och aktivt community → Snabb support & ständiga förbättringar.",
    "• Över 300 000 paket via PyPI – allt från visualisering till maskininlärning.",
    "• Jämfört med MATLAB: färre begränsningar, fler innovationer.",
    "🌍 Global standard\n🚀 Framtidssäkert"
])

# Slide 5: Sammanfattning
add_content_slide("Sammanfattning", [
    "✅ Python är gratis, kraftfullt och flexibelt",
    "✅ Bättre integration med moderna tekniker",
    "✅ Stort community och snabb utveckling",
    "✅ Mer än bara matematik – en hel plattform",
    "👉 För arbetsplatsen betyder det: Lägre kostnader, ökad produktivitet, fler möjligheter."
])

# Spara presentationen
pptx_path = "/mnt/data/Varför_Python_är_bättre_än_MATLAB.pptx"
prs.save(pptx_path)

pptx_path
